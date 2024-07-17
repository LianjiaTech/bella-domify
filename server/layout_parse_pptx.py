# -*- coding: utf-8 -*-
# ===============================================================
#
#    Copyright (C) 2024 Beike, Inc. All Rights Reserved.
#
#    @Create Author : luxu(luxu002@ke.com)
#    @Create Time   : 2024/7/16
#    @Description   : 
#
# ===============================================================
import io
import json
from constants import IMAGE, TEXT, TABLE

from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE

import s3_service
import utils


def analyze_placeholder(shape):
    # 获取占位符的类型和文本内容
    placeholder_type = shape.placeholder_format.type

    if placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:  # 18
        filename = IMAGE + utils.get_random_name()
        image_s3_url = s3_service.s3_upload(filename=filename, data=BytesIO(shape.image.blob))
        return IMAGE, image_s3_url

    else:
        # 文字占位符
        text = shape.text
        return TEXT, text


def layout_parse(file):
    layouts = []

    pptx_stream = io.BytesIO(file)
    pr = Presentation(pptx_stream)
    # 遍历幻灯片中的所有形状
    page_num = 1
    for slide in pr.slides:
        page_num += 1
        for shape in slide.shapes:
            try:
                # placeholder需要识别
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:  # 14
                    shape_type, shape_text = analyze_placeholder(shape)
                    if shape_text:
                        layouts.append(dict(text=shape_text, type=shape_type))

                # 图片
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 13
                    filename = IMAGE + utils.get_random_name()
                    image_s3_url = s3_service.s3_upload(filename=filename, data=BytesIO(shape.image.blob))
                    if image_s3_url:
                        layouts.append(dict(text=image_s3_url, type=IMAGE))

                # 文字
                elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:  # 17
                    if shape.text:
                        layouts.append(dict(text=shape.text, type=TEXT))
                # 表格
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # 19
                    table = shape.table
                    table_text = ""
                    for row in table.rows:
                        for cell in row.cells:
                            table_text = " | ".join([table_text, cell.text])
                    if table_text:
                        layouts.append(dict(text=table_text, type=TABLE))
            except Exception as e:
                print(f"处理元素时出错，type: {shape.shape_type}，errmsg：{str(e)}")
                continue

    result = {"data": {"layouts": layouts}}
    return result


if __name__ == "__main__":

    file_name = '/Users/lucio/code/others/工作内容/工作内容2024/0709多文件类型解析代码调研/demo.pptx'

    # 读取本地文件
    try:
        with open(file_name, 'rb') as file:
            buf_data = file.read()
    except Exception as e:
        print(f"读取文件失败: {e}")

    print(json.dumps(layout_parse(buf_data), ensure_ascii=False))
