# -*- coding: utf-8 -*-
# ===============================================================
#
#    Copyright (C) 2024 Beike, Inc. All Rights Reserved.
#
#    @Create Author : luxu(luxu002@ke.com)
#    @Create Time   : 2024/7/16
#    @Description   : 
#
# ===============================================================
import io
import json

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE

from services.SimpleBlock import SimpleBlock
from services.constants import TEXT, TABLE, IMAGE
from services.layout_parse_utils import get_s3_links_for_simple_block_batch


def layout_parse(file):
    simple_block_list = []

    pptx_stream = io.BytesIO(file)
    pr = Presentation(pptx_stream)
    # 遍历幻灯片中的所有形状
    for slide in pr.slides:
        for shape in slide.shapes:
            try:
                # placeholder需要识别
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:  # 14
                    placeholder_type = shape.placeholder_format.type
                    if placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:  # 18
                        simple_block_list.append(SimpleBlock(type=IMAGE, image_bytes=shape.image.blob))
                    else:
                        # 文字占位符
                        simple_block_list.append(SimpleBlock(type=TEXT, text=shape.text))
                # 图片
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 13
                    simple_block_list.append(SimpleBlock(type=IMAGE, image_bytes=shape.image.blob))
                # 文字
                elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:  # 17
                    if shape.text:
                        simple_block_list.append(SimpleBlock(type=TEXT, text=shape.text))
                # 表格
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # 19
                    table = shape.table
                    table_text = ""
                    for row in table.rows:
                        for cell in row.cells:
                            table_text = " | ".join([table_text, cell.text])
                    if table_text:
                        simple_block_list.append(SimpleBlock(type=TABLE, text=table_text))

            except Exception as e:
                print(f"处理元素时出错，type: {shape.shape_type}，errmsg：{str(e)}")
                continue

    # SimpleBlock的list批量获取S3链接，并返回目标结构
    result = get_s3_links_for_simple_block_batch(simple_block_list)
    return result


if __name__ == "__main__":

    file_name = 'demo.pptx'

    # 读取本地文件
    try:
        with open(file_name, 'rb') as file:
            buf_data = file.read()
    except Exception as e:
        print(f"读取文件失败: {e}")

    print(json.dumps(layout_parse(buf_data), ensure_ascii=False))
